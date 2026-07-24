"""Document parsers — produce paragraph-level RawSegments.

PDF:  PyMuPDF 'dict' mode for block-level layout + font-size heading detection.
DOCX: python-docx style names ('Heading 1', 'Title', ...) drive section tracking.
MD:   ATX-style headings (#, ##, ###) anchor sections.
TXT:  blank-line-separated paragraphs.
XLSX: one segment per non-empty row, section_heading = sheet name.
PPTX: one segment per text frame, section_heading = slide title (or "Slide N").
HTML: <h1..h6> become section headings; nav/script/style/footer are dropped.
CSV:  one row per segment, headers prepended on the first row of each chunk.
"""
from __future__ import annotations

import csv
import io
import logging
import re
from collections import Counter
from collections.abc import Iterator
from io import BytesIO

import fitz  # PyMuPDF
from docx import Document as DocxDocument

from .types import RawSegment

log = logging.getLogger(__name__)


class ParseError(Exception):
    """Raised when a document cannot be parsed."""


class EmptyDocumentError(ParseError):
    """No extractable text — likely a scanned/image PDF or empty file."""


# Heading detection thresholds. Tuned for typical business documents.
_HEADING_FONT_RATIO = 1.18    # block's max font > body_size * ratio → heading
_HEADING_MAX_CHARS = 120      # bold/short blocks above this length are paragraphs, not titles
_BODY_SIZE_SAMPLE_PAGES = 6


def parse_document(file_bytes: bytes, file_type: str) -> list[RawSegment]:
    """Dispatch to the right parser and validate the output."""
    file_type = file_type.lower()
    if file_type == "pdf":
        segments = list(_parse_pdf(file_bytes))
    elif file_type == "docx":
        segments = list(_parse_docx(file_bytes))
    elif file_type == "md":
        segments = list(_parse_markdown(_decode_bytes(file_bytes)))
    elif file_type == "txt":
        segments = list(_parse_text(_decode_bytes(file_bytes)))
    elif file_type == "xlsx":
        segments = list(_parse_xlsx(file_bytes))
    elif file_type == "pptx":
        segments = list(_parse_pptx(file_bytes))
    elif file_type == "html":
        segments = list(_parse_html(file_bytes))
    elif file_type == "csv":
        segments = list(_parse_csv(file_bytes))
    elif file_type in ("vtt", "teams_transcript", "transcript"):
        segments = list(_parse_meeting_transcript(_decode_bytes(file_bytes), file_type=file_type))
    else:
        raise ParseError(f"Unsupported file type: {file_type}")

    if not segments or not any(s.content.strip() for s in segments):
        raise EmptyDocumentError(
            "No extractable text found."
        )
    return segments


# ── PDF ───────────────────────────────────────────────────────────────────────

def _parse_pdf(file_bytes: bytes) -> Iterator[RawSegment]:
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
    except Exception as exc:
        raise ParseError(f"Failed to open PDF: {exc}") from exc

    # Password-protected PDFs (very common for government IDs like Aadhar/PAN)
    # surface as opaque ValueErrors deep inside PyMuPDF's text extraction.
    # Catch up front so the user sees an actionable message instead.
    if getattr(doc, "needs_pass", False) or getattr(doc, "is_encrypted", False):
        try:
            doc.close()
        except Exception:
            pass
        raise ParseError(
            "This PDF is password-protected or encrypted. "
            "Please upload an unprotected copy."
        )

    try:
        body_size = _detect_body_font_size(doc)
        current_heading: str | None = None

        for page in doc:
            page_dict = page.get_text("dict")
            for block in page_dict.get("blocks", []):
                if block.get("type") != 0:  # 0 = text, 1 = image
                    continue

                block_text, max_size, any_bold, line_count = _flatten_block(block)
                if not block_text:
                    continue

                if _is_heading(block_text, max_size, any_bold, line_count, body_size):
                    current_heading = block_text
                    continue

                yield RawSegment(
                    content=block_text,
                    page_number=page.number + 1,
                    section_heading=current_heading,
                    metadata={"font_size": round(max_size, 1)} if max_size else {},
                )
    finally:
        doc.close()


def _flatten_block(block: dict) -> tuple[str, float, bool, int]:
    lines: list[str] = []
    max_size = 0.0
    any_bold = False
    for line in block.get("lines", []):
        spans = line.get("spans", [])
        if not spans:
            continue
        line_text = "".join(s.get("text", "") for s in spans).strip()
        if line_text:
            lines.append(line_text)
        for s in spans:
            size = s.get("size", 0) or 0
            if size > max_size:
                max_size = size
            font_name = s.get("font") or ""
            if "Bold" in font_name or "Black" in font_name:
                any_bold = True
    return " ".join(lines).strip(), max_size, any_bold, len(lines)


def _is_heading(text: str, max_size: float, any_bold: bool, line_count: int, body_size: float) -> bool:
    if len(text) > _HEADING_MAX_CHARS:
        return False
    # Trailing-period paragraphs almost never headings, even if large.
    if text.endswith(".") and len(text) > 60:
        return False
    if body_size and max_size > body_size * _HEADING_FONT_RATIO:
        return True
    if any_bold and line_count <= 2 and len(text) < 80:
        return True
    return False


def _detect_body_font_size(doc) -> float:
    """Body font = the size with the most total character count in the first N pages."""
    sizes: Counter[float] = Counter()
    for page in doc[: min(_BODY_SIZE_SAMPLE_PAGES, len(doc))]:
        for block in page.get_text("dict").get("blocks", []):
            if block.get("type") != 0:
                continue
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    size = round(span.get("size", 0) or 0, 1)
                    text_len = len(span.get("text", ""))
                    if size > 0 and text_len > 0:
                        sizes[size] += text_len
    if not sizes:
        return 11.0
    return sizes.most_common(1)[0][0]


# ── DOCX ──────────────────────────────────────────────────────────────────────

def _parse_docx(file_bytes: bytes) -> Iterator[RawSegment]:
    try:
        doc = DocxDocument(BytesIO(file_bytes))
    except Exception as exc:
        raise ParseError(f"Failed to open DOCX: {exc}") from exc

    current_heading: str | None = None

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style_name = (para.style.name or "").lower() if para.style else ""
        if style_name.startswith("heading") or style_name == "title" or style_name == "subtitle":
            current_heading = text
            continue
        yield RawSegment(content=text, section_heading=current_heading)

    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if not cells:
                continue
            yield RawSegment(
                content=" | ".join(cells),
                section_heading=current_heading,
                metadata={"source": "table"},
            )


# ── Markdown ──────────────────────────────────────────────────────────────────

_MD_HEADING_RE = re.compile(r"^(#{1,6})\s+(.+?)\s*#*\s*$")


def _parse_markdown(text: str) -> Iterator[RawSegment]:
    current_heading: str | None = None
    buffer: list[str] = []

    def drain() -> RawSegment | None:
        joined = "\n".join(buffer).strip()
        buffer.clear()
        if not joined:
            return None
        return RawSegment(content=joined, section_heading=current_heading)

    for line in text.splitlines():
        m = _MD_HEADING_RE.match(line)
        if m:
            seg = drain()
            if seg:
                yield seg
            current_heading = m.group(2).strip()
            continue
        if not line.strip():
            seg = drain()
            if seg:
                yield seg
            continue
        buffer.append(line)

    seg = drain()
    if seg:
        yield seg


# ── Plain text ────────────────────────────────────────────────────────────────

def _parse_text(text: str) -> Iterator[RawSegment]:
    for para in re.split(r"\n\s*\n", text):
        para = para.strip()
        if para:
            yield RawSegment(content=para)


# ── Meeting transcripts (Zoom VTT / Teams JSON) ───────────────────────────────

def _parse_meeting_transcript(text: str, *, file_type: str) -> Iterator[RawSegment]:
    """Project a transcript to `Speaker: text` paragraph segments.

    We chunk by speaker turn rather than blank-line paragraphs so retrieval
    surfaces utterance-level matches ("what did Alice say about X") with
    the speaker label intact. Empty transcripts raise EmptyDocumentError
    via the outer parse_document validation, same as scanned PDFs.
    """
    from app.services.parsers.meeting_transcript import (
        detect_transcript_format,
        parse_transcript,
    )

    # Day 13 production-hardening: when a user uploads a .json marked as
    # teams_transcript, distinguish "real Teams export with no recognised
    # speakers" (rare, EmptyDocumentError downstream is fine) from "this
    # isn't a Teams transcript at all" (common — they uploaded the wrong
    # JSON). The second case gets a friendlier ParseError that tells them
    # what shape we expect.
    if file_type == "teams_transcript":
        detected = detect_transcript_format(file_type=None, content=text)
        if detected != "teams_json":
            raise ParseError(
                "This JSON doesn't look like a Microsoft Teams transcript. "
                "Teams exports have a top-level `recognizedPhrases` or "
                "`entries` array. Export your transcript from Teams (Meeting "
                "details → Recordings & transcripts → Download)."
            )

    parsed = parse_transcript(file_type=file_type, content=text)
    if parsed.is_empty():
        return
    # Group consecutive utterances by the same speaker into a single segment.
    # Improves retrieval: a single contiguous answer reads better than 8
    # fragmented "Alice: …" lines, and the chunker re-splits if needed.
    cur_speaker: str | None = None
    cur_lines: list[str] = []
    for u in parsed.utterances:
        if u.speaker == cur_speaker:
            cur_lines.append(u.text)
            continue
        if cur_speaker and cur_lines:
            yield RawSegment(
                content=f"{cur_speaker}: " + " ".join(cur_lines).strip(),
                section_heading=cur_speaker,
            )
        cur_speaker = u.speaker
        cur_lines = [u.text]
    if cur_speaker and cur_lines:
        yield RawSegment(
            content=f"{cur_speaker}: " + " ".join(cur_lines).strip(),
            section_heading=cur_speaker,
        )


# ── XLSX ──────────────────────────────────────────────────────────────────────

# Cap per-sheet row scan to keep parsing memory bounded on hostile inputs.
# 100k rows × ~32 cells × ~16 chars ≈ 50 MB of text, which is comfortably
# within our memory budget and well past any realistic business spreadsheet.
_XLSX_MAX_ROWS_PER_SHEET = 100_000


def _parse_xlsx(file_bytes: bytes) -> Iterator[RawSegment]:
    try:
        import openpyxl  # local import: keeps cold-start fast for PDF-only orgs
    except ImportError as exc:
        raise ParseError("openpyxl is not installed") from exc

    try:
        wb = openpyxl.load_workbook(BytesIO(file_bytes), read_only=True, data_only=True)
    except Exception as exc:
        raise ParseError(f"Failed to open XLSX: {exc}") from exc

    try:
        for sheet in wb.worksheets:
            sheet_title = (sheet.title or "Sheet").strip()
            header: list[str] = []
            for row_idx, row in enumerate(sheet.iter_rows(values_only=True)):
                if row_idx >= _XLSX_MAX_ROWS_PER_SHEET:
                    log.warning("xlsx: %s truncated at %d rows", sheet_title, _XLSX_MAX_ROWS_PER_SHEET)
                    break
                cells = [_xlsx_cell_to_text(c) for c in row]
                if not any(cells):
                    continue
                if not header and any(cells):
                    # Treat the first non-empty row as the header so later rows
                    # are emitted as "col: value" lines — much more retrievable
                    # than a bare CSV-style "a | b | c".
                    header = cells
                    yield RawSegment(
                        content=" | ".join(c for c in cells if c),
                        section_heading=sheet_title,
                        metadata={"row": row_idx + 1, "source": "header"},
                    )
                    continue
                if header:
                    pairs = [
                        f"{header[i]}: {cells[i]}"
                        for i in range(min(len(header), len(cells)))
                        if cells[i]
                    ]
                    text = ", ".join(pairs) if pairs else " | ".join(c for c in cells if c)
                else:
                    text = " | ".join(c for c in cells if c)
                if text.strip():
                    yield RawSegment(
                        content=text,
                        section_heading=sheet_title,
                        metadata={"row": row_idx + 1},
                    )
    finally:
        wb.close()


def _xlsx_cell_to_text(value) -> str:
    if value is None:
        return ""
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value).strip()


# ── PPTX ──────────────────────────────────────────────────────────────────────

def _parse_pptx(file_bytes: bytes) -> Iterator[RawSegment]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ParseError("python-pptx is not installed") from exc

    try:
        prs = Presentation(BytesIO(file_bytes))
    except Exception as exc:
        raise ParseError(f"Failed to open PPTX: {exc}") from exc

    for slide_idx, slide in enumerate(prs.slides, start=1):
        title = _pptx_slide_title(slide) or f"Slide {slide_idx}"

        for shape in slide.shapes:
            text = _pptx_shape_text(shape)
            if text and text != title:
                yield RawSegment(
                    content=text,
                    section_heading=title,
                    metadata={"slide": slide_idx},
                )

        notes = _pptx_notes(slide)
        if notes:
            yield RawSegment(
                content=notes,
                section_heading=title,
                metadata={"slide": slide_idx, "source": "notes"},
            )


def _pptx_slide_title(slide) -> str | None:
    try:
        if slide.shapes.title and slide.shapes.title.text:
            return slide.shapes.title.text.strip() or None
    except Exception:
        return None
    return None


def _pptx_shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    parts: list[str] = []
    for para in shape.text_frame.paragraphs:
        line = "".join(run.text for run in para.runs).strip()
        if line:
            parts.append(line)
    return "\n".join(parts).strip()


def _pptx_notes(slide) -> str:
    try:
        if not getattr(slide, "has_notes_slide", False) or not slide.has_notes_slide:
            return ""
        nf = slide.notes_slide.notes_text_frame
        return (nf.text or "").strip()
    except Exception:
        return ""


# ── HTML ──────────────────────────────────────────────────────────────────────

_HTML_DROP_TAGS = ("script", "style", "noscript", "template", "nav", "footer", "aside", "header")


def _parse_html(file_bytes: bytes) -> Iterator[RawSegment]:
    try:
        from bs4 import BeautifulSoup
    except ImportError as exc:
        raise ParseError("beautifulsoup4 is not installed") from exc

    text = _decode_bytes(file_bytes)
    try:
        soup = BeautifulSoup(text, "lxml")
    except Exception:
        # lxml may be missing in dev — fall back to the stdlib parser so HTML
        # still ingests, just with slightly worse heading detection.
        soup = BeautifulSoup(text, "html.parser")

    for tag in soup(_HTML_DROP_TAGS):
        tag.decompose()

    current_heading: str | None = None
    root = soup.body or soup
    for element in root.descendants:
        name = getattr(element, "name", None)
        if name is None:
            continue
        if name in ("h1", "h2", "h3", "h4", "h5", "h6"):
            heading = element.get_text(" ", strip=True)
            if heading:
                current_heading = heading[:200]
            continue
        if name in ("p", "li", "blockquote", "td", "th", "dd", "dt", "pre"):
            body = element.get_text(" ", strip=True)
            if body and len(body) > 1:
                yield RawSegment(content=body, section_heading=current_heading)


# ── CSV ───────────────────────────────────────────────────────────────────────

_CSV_MAX_ROWS = 100_000


def _parse_csv(file_bytes: bytes) -> Iterator[RawSegment]:
    text = _decode_bytes(file_bytes)
    sample = text[:8192]
    try:
        dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
    except csv.Error:
        dialect = csv.excel

    reader = csv.reader(io.StringIO(text), dialect=dialect)
    header: list[str] = []
    for row_idx, row in enumerate(reader):
        if row_idx >= _CSV_MAX_ROWS:
            log.warning("csv: truncated at %d rows", _CSV_MAX_ROWS)
            break
        cells = [c.strip() for c in row]
        if not any(cells):
            continue
        if not header:
            header = cells
            yield RawSegment(
                content=" | ".join(c for c in cells if c),
                metadata={"row": row_idx + 1, "source": "header"},
            )
            continue
        pairs = [
            f"{header[i]}: {cells[i]}"
            for i in range(min(len(header), len(cells)))
            if cells[i]
        ]
        if pairs:
            yield RawSegment(content=", ".join(pairs), metadata={"row": row_idx + 1})


# ── Bytes → str with fallback ─────────────────────────────────────────────────

def _decode_bytes(file_bytes: bytes) -> str:
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return file_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("utf-8", errors="replace")
